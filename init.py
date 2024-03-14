from pptx import Presentation
from dataclasses import dataclass
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor
import datetime
import pyautogui
import time
import csv
import numpy as np
import os
import pandas as pd
import matplotlib.pyplot as plt

@dataclass
class Account:
    name: str
    chrome_profile_path: str
    url_results: str
    url_people: str
    url_content_summary: str
    download_dir: str = 'C:/Users/hoorm/Downloads'
    chrome_data_path: str = "C:/Users/hoorm/AppData/Local/Google/Chrome/User Data/"

accounts = [
    Account("BDMH", 
            "Profile 4", 
            "https://business.facebook.com/latest/insights/results?business_id=693823586008930&asset_id=188285954362840&ad_account_id=120200701245850619&entity_type=FB_PAGE", 
            "https://business.facebook.com/latest/insights/people?business_id=693823586008930&asset_id=188285954362840&ad_account_id=120200701245850619&entity_type=FB_PAGE", 
            "https://business.facebook.com/latest/insights/content_summary?business_id=693823586008930&asset_id=188285954362840&ad_account_id=120200701245850619&entity_type=FB_PAGE"),
    Account("HAKC",
            "Default",
            "https://business.facebook.com/latest/insights/results/?entity_type=FB_PAGE&ad_account_id=120205803738370462&asset_id=192596830610881&business_id=1499208007321158",
            "https://business.facebook.com/latest/insights/people?entity_type=FB_PAGE&ad_account_id=120205803738370462&asset_id=192596830610881&business_id=1499208007321158",
            "https://business.facebook.com/latest/insights/content_summary?entity_type=FB_PAGE&ad_account_id=120205803738370462&asset_id=192596830610881&business_id=1499208007321158"),
]

for account in accounts:
    # Get a list of all files in the download directory before the download
    before = os.listdir(account.download_dir)

    # Create Chrome options and set the profile path
    chrome_options = Options()
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument(f'--user-data-dir={account.chrome_data_path}')
    chrome_options.add_argument(f'--profile-directory={account.chrome_profile_path}')
    chrome_options.add_argument("--disable-dev-shm-usage")
    chrome_options.add_experimental_option('excludeSwitches', ['enable-logging'])

    # Create a new instance of the Chrome driver with the specified options
    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=chrome_options)

    # Navigate to the specified URL
    driver.get(account.url_results)
    print("Navigating to the URL, resuming in 8 seconds...")
    time.sleep(8)

    matches = pyautogui.locateAllOnScreen('img/Export.png', confidence=0.8)
    counter = 0
    for match in matches:
        if counter == 1:
            print(match)
            pyautogui.click(match)
            time.sleep(1)
            download = pyautogui.locateOnScreen('img/csv.png', confidence=0.95)
            pyautogui.click(download)
            time.sleep(1)
        counter += 1

    print("Downloading the file, resuming in 5 seconds...")    
    time.sleep(5)

    # Get a list of all files in the download directory after the download
    after = os.listdir(account.download_dir)

    downloaded_files = [f for f in after if f not in before]
    # Get the path of the downloaded file
    downloaded_file_path = os.path.join(account.download_dir, downloaded_files[0])
    # Read the CSV file into a pandas DataFrame

    with open(downloaded_file_path, 'r', encoding='UTF-16') as f:
        lines = f.readlines()

    lines_with_comma = [line for line in lines if ',' in line]

    with open(downloaded_file_path, 'w', encoding='UTF-16') as f:
        f.writelines(lines_with_comma)

    def read_multiple_tables(file):
        with open(file, 'r', encoding='UTF-16') as f:
            lines = f.readlines()

        tables = []
        table = []
        for line in lines[1:]:  # skip the first line ("sep=,")
            if 'Date' in line:  # this is a header line
                if table:  # if there is already a table, add it to the list of tables
                    tables.append(pd.DataFrame(table[1:], columns=table[0]))
                table = [line.strip().split(',')]  # start a new table
            else:
                table.append(line.strip().split(','))  # add a row to the current table
        if table:  # add the last table to the list of tables
            tables.append(pd.DataFrame(table[1:], columns=table[0]))

        return tables

    tables = read_multiple_tables(downloaded_file_path)

    facebook_likes_sum = 0
    instagram_followers_sum = 0

    for table in tables:
        table.columns = table.columns.str.replace('"', '')
        for col in table.columns:
            table[col] = table[col].str.replace('"', '')
            
        if ('Facebook Page likes' in table.columns.values[1]):
            facebook_likes_sum = table["Facebook Page likes"].astype(int).sum()
            print(table.columns, facebook_likes_sum, instagram_followers_sum)

        elif ("Instagram followers" in table.columns.values[1]):
            instagram_followers_sum = table["Instagram followers"].astype(int).sum()
            print(table.columns, facebook_likes_sum, instagram_followers_sum)


    # Delete the downloaded file
    os.remove(downloaded_file_path)
    # print(visitsData)

    # Navigate to the specified URL
    driver.get(account.url_people)
    print("Navigating to the URL, resuming in 8 seconds...")
    time.sleep(8)

    # Wait for the element to be clickable
    match = pyautogui.locateOnScreen('img/Export.png', confidence=0.8)
    pyautogui.click(match)
    print("Clicking the Export button, resuming in 1 second...")
    time.sleep(1)
    download = pyautogui.locateOnScreen('img/csv.png', confidence=0.95)
    pyautogui.click(download)
    print("Clicking the CSV button, resuming in 1 second...")
    time.sleep(1)

    # Get a list of all files in the download directory after the download
    after = os.listdir(account.download_dir)

    downloaded_files = [f for f in after if f not in before]
    # Get the path of the downloaded file
    downloaded_file_path = os.path.join(account.download_dir, downloaded_files[0])
    # Read the CSV file into a pandas DataFrame, skipping the problematic line

    with open(downloaded_file_path, 'r', encoding='UTF-16') as f:
        lines = f.readlines()

    pre_lines = []
    new_lines = []

    for line in lines:
        newline = ""
        qopen = False
        for char in line:
            if char == '"':
                qopen = not qopen
            elif qopen and not char == ',':
                newline += char
            elif not qopen:
                newline += char
        newline = newline.replace('"', '')
        pre_lines.append(newline)

    i = 0
    while i < len(pre_lines):
        line = pre_lines[i].strip()  # remove leading and trailing whitespaces
        if line:  # if the line is not empty
            if ',' not in line and not line.isdigit():  # if the line does not have a comma
                if i + 1 < len(pre_lines):  # if there is a next line
                    new_lines.append("*" + line + pre_lines[i + 1])  # append the current line and the next line
                    i += 1  # skip the next line
            else:
                new_lines.append(line)
                new_lines.append('\n')  # add a newline character
        i += 1

    with open(downloaded_file_path, 'w', encoding='UTF-16') as f:
        f.writelines(new_lines)




    def read_multiple_tables(file):
        with open(file, 'r', encoding='UTF-16') as f:
            lines = f.readlines()

        tables = []
        table = []
        for line in lines[1:]:  # skip the first line ("sep=,")
            if '*' in line:  # this is a header line
                if table:  # if there is already a table, add it to the list of tables
                    tables.append(pd.DataFrame(table[1:], columns=table[0]))
                table = [line.strip().split(',')]  # start a new table
            else:
                table.append(line.strip().split(','))  # add a row to the current table
        if table:  # add the last table to the list of tables
            tables.append(pd.DataFrame(table[1:], columns=table[0]))

        return tables

    tables = read_multiple_tables(downloaded_file_path)

    totalFollowersIG = 0
    totalFollowersFB = 0

    topCityIG = "N/A"
    topCityFB = "N/A"

    topCountryIG = "N/A"
    topCountryFB = "N/A"

    topAgeIG = "N/A"
    topAgeFB = "N/A"

    for table in tables:
        if "*Facebook followersFB_PAGEFOLLOWUNIQUE_USERS" in table.columns.values[0]:
            totalFollowersFB = table["*Facebook followersFB_PAGEFOLLOWUNIQUE_USERS"][0]
            print(totalFollowersFB)
        
        if "*Instagram followersIG_ACCOUNTFOLLOWUNIQUE_USERS" in table.columns.values[0]:
            totalFollowersIG = table["*Instagram followersIG_ACCOUNTFOLLOWUNIQUE_USERS"][0]
            print(totalFollowersIG)
        
        if "*Facebook followers by gender and age" in table.columns.values[0]:
            if not table.empty:
                table['Women'] = table['Women'].str.rstrip('%').astype('float')
                table['Men'] = table['Men'].str.rstrip('%').astype('float')

                max_women = table['Women'].idxmax()
                max_men = table['Men'].idxmax()

                if table.loc[max_women, 'Women'] > table.loc[max_men, 'Men']:
                    topAgeFB = f"Women, {table.loc[max_women, '*Facebook followers by gender and ageAge']}, with {table.loc[max_women, 'Women']}%"
                else:
                    topAgeFB = f"Men, {table.loc[max_men, '*Facebook followers by gender and ageAge']}, with {table.loc[max_men, 'Men']}%"
            print(topAgeFB)
        
        if "*Instagram followers by gender and age" in table.columns.values[0]:
            if not table.empty:
                table['Women'] = table['Women'].str.rstrip('%').astype('float')
                table['Men'] = table['Men'].str.rstrip('%').astype('float')

                max_women = table['Women'].idxmax()
                max_men = table['Men'].idxmax()

                if table.loc[max_women, 'Women'] > table.loc[max_men, 'Men']:
                    topAgeIG = f"Women, {table.loc[max_women, '*Instagram followers by gender and ageAge']}, with {table.loc[max_women, 'Women']}%"
                else:
                    topAgeIG = f"Men, {table.loc[max_men, '*Instagram followers by gender and ageAge']}, with {table.loc[max_men, 'Men']}%"
            print(topAgeIG)

        if "*Facebook followers by top countries" in table.columns.values[0]:
            if not table.empty:
                table['Value'] = table['Value'].str.rstrip('%').astype('float')
                max_country = table['Value'].idxmax()
                topCountryFB = f"{table.loc[max_country, '*Facebook followers by top countriesTop countries']}, with {table.loc[max_country, 'Value']}%"
            print(topCountryFB)

        if "*Instagram followers by top countries" in table.columns.values[0]:
            if not table.empty:
                table['Value'] = table['Value'].str.rstrip('%').astype('float')
                max_country = table['Value'].idxmax()
                topCountryIG = f"{table.loc[max_country, '*Instagram followers by top countriesTop countries']}, with {table.loc[max_country, 'Value']}%"
            print(topCountryIG)


    # Delete the downloaded file
    os.remove(downloaded_file_path)

    # Navigate to the specified URL
    driver.get(account.url_content_summary)
    print("Navigating to the URL, resuming in 8 seconds...")
    time.sleep(8)


    # Take a screenshot of a specific region
    x = 980  # X-coordinate of the top-left corner of the region
    y = 390  # Y-coordinate of the top-left corner of the region
    width = 930  # Width of the region
    height = 430  # Height of the region
    screenshot = pyautogui.screenshot(region=(x, y, width, height))
    # Save the screenshot as a PNG file
    screenshot.save('FB_Reach.png')

    # Take a screenshot of a specific region
    x = 980  # X-coordinate of the top-left corner of the region
    y = 810  # Y-coordinate of the top-left corner of the region
    width = 930  # Width of the region
    height = 430  # Height of the region
    screenshot = pyautogui.screenshot(region=(x, y, width, height))
    # Save the screenshot as a PNG file
    screenshot.save('FB_Engagement.png')

    download = pyautogui.locateOnScreen('img/instagramposts.png', confidence=0.95)
    pyautogui.click(download)
    print("Navigating to Instagram page, resuming in 5 seconds...")
    time.sleep(5)

    # Take a screenshot of a specific region
    x = 980  # X-coordinate of the top-left corner of the region
    y = 390  # Y-coordinate of the top-left corner of the region
    width = 930  # Width of the region
    height = 400  # Height of the region
    screenshot = pyautogui.screenshot(region=(x, y, width, height))
    # Save the screenshot as a PNG file
    screenshot.save('IG_Reach.png')

    # Take a screenshot of a specific region
    x = 980  # X-coordinate of the top-left corner of the region
    y = 790  # Y-coordinate of the top-left corner of the region
    width = 930  # Width of the region
    height = 400  # Height of the region
    screenshot = pyautogui.screenshot(region=(x, y, width, height))
    # Save the screenshot as a PNG file
    screenshot.save('IG_Engagement.png')

    # Open the PowerPoint file
    pptx_file = 'ppt/SocialMedia_Report.pptx'
    presentation = Presentation(pptx_file)

    # Access the slides in the presentation
    slides = presentation.slides

    # Insert a new slide with heading "Instagram Insights"
    slide = slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Instagram Insights"
    content = slide.placeholders[1]
    content.text = f"Instagram Visits: {instagram_followers_sum}\nTotal Followers: {totalFollowersIG}\nTop Visitors: {topAgeIG}\nTop Country: {topCountryIG}"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(37, 37, 37)  # Set background color to #252525
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to #FFFFFF


    # Insert a slide with the photo ./IG_Reach.png
    slide = slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture("./IG_Reach.png", left=0, top=0)

    # Insert a slide with the photo ./IG_Engagement.png
    slide = slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture("./IG_Engagement.png", left=0, top=0)

    # Insert a new slide with heading "Facebook Insights"
    slide = slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Facebook Insights"
    content = slide.placeholders[1]
    content.text = f"Facebook Visits: {facebook_likes_sum}\nTotal Followers: {totalFollowersFB}\nTop Visitors: {topAgeFB}\nTop Country: {topCountryFB}"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(37, 37, 37)  # Set background color to #252525
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to #FFFFFF

    # Insert a slide with the photo ./IG_Reach.png
    slide = slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture("./FB_Reach.png", left=0, top=0)

    # Insert a slide with the photo ./IG_Engagement.png
    slide = slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture("./FB_Engagement.png", left=0, top=0)

    # Insert a slide with the photo closing slide
    slide = slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture("./img/closing.png", left=0, top=0, width=presentation.slide_width, height=presentation.slide_height)

    # Save the modified PowerPoint file
    current_date = datetime.datetime.now().strftime("%Y-%m-%d")
    modified_pptx_file = f'report_{account.name}_{current_date}.pptx'
    presentation.save(modified_pptx_file)

    # Close the browser
    driver.quit()